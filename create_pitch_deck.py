#!/usr/bin/env python3
"""
快速生成路演PPT的脚本
使用前请先安装: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_pitch_deck():
    """创建5分钟路演PPT"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 宽屏
    prs.slide_height = Inches(7.5)
    
    # 定义颜色
    PRIMARY_COLOR = RGBColor(30, 58, 95)  # #1e3a5f 深蓝
    ACCENT_COLOR = RGBColor(59, 130, 246)  # #3b82f6 亮蓝
    SUCCESS_COLOR = RGBColor(16, 185, 129)  # #10b981 绿色
    
    # ========== 第1页: 封面 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    
    title = slide.shapes.title
    title.text = "EvoAgent-AFP"
    title.text_frame.paragraphs[0].font.size = Pt(60)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    subtitle = slide.placeholders[1]
    subtitle.text = "自进化智能体驱动的抗冻蛋白精准设计平台\n\nAI+X 创新应用大赛决赛路演\n汇报人: 郭群飞 | BGI华大基因华中分院"
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[1].font.size = Pt(20)
    
    # ========== 第2页: 痛点与机遇 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "为什么需要AI+抗冻蛋白?"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5))
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "🔴 行业痛点"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(220, 38, 38)
    
    p = tf.add_paragraph()
    p.text = "\n• 耗时漫长: 实验筛选 3-6个月/轮"
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "• 成本高昂: 单次突变筛选 ¥50,000+"
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "• 成功率低: 随机突变有效率 < 5%"
    p.font.size = Pt(20)
    
    # 右侧 - AI带来的变革
    content2 = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(6), Inches(5))
    tf2 = content2.text_frame
    
    p = tf2.add_paragraph()
    p.text = "✅ AI带来的变革"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_COLOR
    
    p = tf2.add_paragraph()
    p.text = "\n✓ 设计时间: 数月 → 数小时"
    p.font.size = Pt(20)
    
    p = tf2.add_paragraph()
    p.text = "✓ 研发成本: 降低 90%"
    p.font.size = Pt(20)
    
    p = tf2.add_paragraph()
    p.text = "✓ 预测精度: 提升至 87%"
    p.font.size = Pt(20)
    
    p = tf2.add_paragraph()
    p.text = "✓ 知识沉淀: 可复用、可传承"
    p.font.size = Pt(20)
    
    # ========== 第3页: 解决方案 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "EvoAgent-AFP: 一站式智能设计平台"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "核心能力:"
    p.font.size = Pt(24)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "\n1️⃣ 多模型AI引擎: 集成7大主流LLM(Kimi/DeepSeek/GPT-5.5/Qwen/GLM/MiniMax)"
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "2️⃣ 智能工作流: 序列分析 → 智能突变 → 活性评估 → 冰结合模拟 → 迭代优化"
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "3️⃣ 7个AFP专用工具: afp_sequence_analyze | afp_mutate_sequence | afp_evaluate_mutation | ..."
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "4️⃣ 批量处理: 单次500+序列,实时进度追踪,CSV/JSON导出"
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "5️⃣ 自进化架构: 持续学习积累经验,越用越聪明"
    p.font.size = Pt(20)
    
    # ========== 第4页: 技术创新 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "四大核心技术突破"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.8))
    tf = content.text_frame
    
    innovations = [
        ("1️⃣ 自进化智能体 🔄", "Agent自动调用工具链完成任务,设计结果反馈到知识库持续学习"),
        ("2️⃣ 多模型协同决策 🤖×7", "同时调用7个顶级AI模型交叉验证,实测准确率提升至87%"),
        ("3️⃣ 领域知识深度融合 🧬", "构建AFP专属知识库,RAG检索增强,冰结合基序规则嵌入"),
        ("4️⃣ 端到端自动化 ⚡", "从序列输入到设计方案全自动,效率提升100倍"),
    ]
    
    for tech_title, tech_desc in innovations:
        p = tf.add_paragraph()
        p.text = f"\n{tech_title}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        
        p = tf.add_paragraph()
        p.text = f"   {tech_desc}"
        p.font.size = Pt(18)
    
    # ========== 第5页: 应用案例 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "实战效果验证"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.8))
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "案例: AFP序列优化设计"
    p.font.size = Pt(24)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "\n输入序列(野生型, 33aa):"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "DTASDAAAAAALTAANAAAAAKLTADNAAAAAAATAR"
    p.font.size = Pt(16)
    p.font.name = "Courier New"
    
    p = tf.add_paragraph()
    p.text = "\nAI分析结果:"
    p.font.size = Pt(20)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "✓ 识别冰结合基序: T-x-T重复模式"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "✓ 推荐突变位点: A12T, N15T, K18A"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "✓ 预期TH活性提升: +150%"
    p.font.size = Pt(18)
    p.font.color.rgb = SUCCESS_COLOR
    
    # 性能对比表格
    p = tf.add_paragraph()
    p.text = "\n性能对比:"
    p.font.size = Pt(20)
    p.font.bold = True
    
    table_data = [
        ["指标", "传统方法", "EvoAgent-AFP", "提升"],
        ["设计时间", "2周", "5分钟", "400倍 ⚡"],
        ["候选方案", "3-5个", "20+个", "4倍 📊"],
        ["成功率", "5%", "35%", "7倍 🎯"],
        ["单次成本", "¥50,000", "¥500", "100倍 💰"],
    ]
    
    rows = len(table_data)
    cols = len(table_data[0])
    left = Inches(1)
    top = Inches(5.5)
    width = Inches(10.5)
    height = Inches(1.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    for i, row in enumerate(table_data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = cell_text
            if i == 0:  # 表头
                cell.text_frame.paragraphs[0].font.size = Pt(14)
                cell.text_frame.paragraphs[0].font.bold = True
            else:
                cell.text_frame.paragraphs[0].font.size = Pt(12)
    
    # ========== 第6页: 市场价值 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "商业前景广阔"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.8))
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "目标市场 (总计¥850亿):"
    p.font.size = Pt(22)
    p.font.bold = True
    
    markets = [
        "🏥 生物医药: ¥500亿 (细胞冷冻+器官保存)",
        "🍦 食品工业: ¥200亿 (冰淇淋+冷链保鲜)",
        "🌾 农业育种: ¥100亿 (抗寒作物+胚胎保存)",
        "🔬 科研服务: ¥50亿 (CRO+定制设计)",
    ]
    
    for market in markets:
        p = tf.add_paragraph()
        p.text = f"• {market}"
        p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "\n商业模式:"
    p.font.size = Pt(20)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "💰 SaaS订阅: ¥5,000-20,000/月 | 按需付费: ¥100-500/次"
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "💰 定制开发: ¥50,000-200,000/项目 | API授权: ¥10,000/月"
    p.font.size = Pt(16)
    
    # ========== 第7页: 发展规划 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "未来路线图"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.8))
    tf = content.text_frame
    
    plans = [
        ("2026 (已完成) ✅", [
            "✓ 平台MVP开发完成",
            "✓ 支持7大AI模型",
            "✓ 实现批量处理功能",
        ]),
        ("2027 (短期) 🎯", [
            "• 积累100+企业用户",
            "• 建立万级AFP序列数据库",
            "• Pre-A轮融资¥1000万",
        ]),
        ("2028-2030 (长期) 🌟", [
            "• 成为蛋白质设计标杆平台",
            "• 服务1000+企业和科研机构",
            "• 年营收突破¥1亿",
        ]),
    ]
    
    for period, items in plans:
        p = tf.add_paragraph()
        p.text = f"\n{period}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
    
    # ========== 第8页: 团队介绍 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "核心团队"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.8))
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "郭群飞 - 创始人&CEO"
    p.font.size = Pt(24)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "🎓 生物学博士,BGI资深研究员 | 💼 10年蛋白质工程经验 | 🏆 SCI论文20+篇"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "\n技术团队 (8人):"
    p.font.size = Pt(20)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• AI算法工程师 × 3"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• 生物信息学家 × 2"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• 全栈工程师 × 2 | 产品经理 × 1"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "\n顾问团队: 蛋白质结构预测专家 | AI药物设计权威 | 生物制药CTO"
    p.font.size = Pt(16)
    
    # ========== 第9页: 致谢 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # 居中大字
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = "感谢聆听!"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    textbox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1))
    tf2 = textbox2.text_frame
    p = tf2.add_paragraph()
    p.text = "让AI赋能生命科学研究 · 加速抗冻蛋白创新应用"
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    
    # 联系方式
    textbox3 = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9), Inches(2))
    tf3 = textbox3.text_frame
    p = tf3.add_paragraph()
    p.text = "📧 guoqunfei@bgigroup.com\n🌐 [演示网站URL]\n\n期待与您合作,共同开创AI+生物新时代!"
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    
    # 保存PPT
    output_path = "/Users/guoqunfei/Desktop/BGI/biodiversity/other/AI-抗冻蛋白/EvoAgent-AFP_路演PPT.pptx"
    prs.save(output_path)
    print(f"✅ PPT已生成: {output_path}")
    print("\n📝 使用说明:")
    print("1. 打开生成的PPT文件")
    print("2. 根据需要调整字体、颜色和布局")
    print("3. 添加图表、图标和截图")
    print("4. 插入公司Logo和产品截图")
    print("5. 排练计时,确保总时长在5分钟内")
    print("\n祝路演成功! 🎉")

if __name__ == "__main__":
    try:
        create_pitch_deck()
    except ImportError:
        print("❌ 请先安装python-pptx: pip install python-pptx")
        print("\n或者直接使用 PITCH_DECK_SCRIPT.md 中的内容手动制作PPT")
