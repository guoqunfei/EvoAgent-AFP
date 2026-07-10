#!/usr/bin/env python3
"""
将项目PPT内容整合到参考模板中
生成符合比赛要求的路演PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import shutil

# 文件路径
TEMPLATE_PATH = '附件2.【参考模板】决赛路演--姓名.pptx'
PROJECT_PPT_PATH = 'EvoAgent-AFP：自进化智能体驱动的抗冻蛋白精准设计平台.pptx'
OUTPUT_PATH = 'EvoAgent-AFP_决赛路演PPT_郭群飞.pptx'

def extract_project_content():
    """提取项目PPT中的所有内容"""
    project_ppt = Presentation(PROJECT_PPT_PATH)
    slides_content = []
    
    for i, slide in enumerate(project_ppt.slides):
        slide_data = {
            'title': '',
            'texts': [],
            'images': []
        }
        
        # 提取标题
        if slide.shapes.title:
            slide_data['title'] = slide.shapes.title.text
        
        # 提取所有文本内容
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape != slide.shapes.title:
                    slide_data['texts'].append({
                        'text': shape.text,
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height
                    })
        
        # 提取图片
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                slide_data['images'].append(shape)
        
        slides_content.append(slide_data)
    
    return slides_content

def create_final_presentation():
    """创建最终的路演PPT"""
    
    # 复制参考模板作为基础
    shutil.copy(TEMPLATE_PATH, OUTPUT_PATH)
    prs = Presentation(OUTPUT_PATH)
    
    print(f"✅ 已加载参考模板: {len(prs.slides)}页")
    
    # 提取项目内容
    project_content = extract_project_content()
    print(f"✅ 已提取项目内容: {len(project_content)}页")
    
    # 映射关系: 项目页面 -> 模板页面
    # 根据评分表和比赛要求,重新组织内容结构
    content_mapping = {
        0: {'template_page': 0, 'title': '封面'},
        1: {'template_page': 1, 'title': '应用场景与痛点'},
        2: {'template_page': 2, 'title': '解决方案'},
        3: {'template_page': 3, 'title': '核心技术'},
        4: {'template_page': 4, 'title': '效果验证'},
        5: {'template_page': 5, 'title': '商业模式'},
        6: {'template_page': 6, 'title': '团队介绍'},
    }
    
    # 更新每一页的内容
    for mapping in content_mapping.values():
        proj_idx = mapping['template_page']
        if proj_idx < len(project_content):
            slide = prs.slides[mapping['template_page']]
            content = project_content[proj_idx]
            
            print(f"\n📝 处理第{mapping['template_page'] + 1}页: {mapping['title']}")
            
            # 删除原有的文本框(保留背景和设计元素)
            shapes_to_remove = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # 只删除文本框,保留其他元素
                    if shape.shape_type == 17:  # Text box
                        shapes_to_remove.append(shape)
            
            # 添加项目内容
            if content['texts']:
                # 如果有标题,添加到标题位置
                if content['title'] and slide.shapes.title:
                    slide.shapes.title.text = content['title']
                
                # 添加文本内容
                for j, text_item in enumerate(content['texts']):
                    try:
                        # 在原有位置添加文本框
                        textbox = slide.shapes.add_textbox(
                            text_item['left'],
                            text_item['top'],
                            text_item['width'],
                            text_item['height']
                        )
                        tf = textbox.text_frame
                        tf.text = text_item['text']
                        
                        # 设置字体样式
                        for paragraph in tf.paragraphs:
                            paragraph.font.size = Pt(18)
                            paragraph.font.name = '微软雅黑'
                    except Exception as e:
                        print(f"   ⚠️  添加文本项{j}失败: {e}")
    
    # 保存
    prs.save(OUTPUT_PATH)
    print(f"\n{'=' * 80}")
    print(f"✅ 路演PPT生成成功!")
    print(f"{'=' * 80}")
    print(f"📁 文件位置: {OUTPUT_PATH}")
    print(f"📊 总页数: {len(prs.slides)}页")
    print(f"\n📋 PPT结构:")
    
    for i, slide in enumerate(prs.slides, 1):
        title = slide.shapes.title.text if slide.shapes.title else "无标题"
        print(f"   {i}. {title}")
    
    print(f"\n🎯 下一步:")
    print(f"   1. 打开PPT检查内容和排版")
    print(f"   2. 调整字体大小和颜色")
    print(f"   3. 添加必要的图表和图片")
    print(f"   4. 排练演讲(目标: 5分钟)")

if __name__ == '__main__':
    try:
        create_final_presentation()
    except Exception as e:
        print(f"❌ 错误: {e}")
        import traceback
        traceback.print_exc()
