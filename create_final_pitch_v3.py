#!/usr/bin/env python3
"""
终极版路演PPT生成器
完全基于原始项目PPT的内容和设计风格,按照参考模板的结构重组
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import shutil

# 文件路径
TEMPLATE_PATH = '附件2.【参考模板】决赛路演--姓名.pptx'
PROJECT_PPT_PATH = 'EvoAgent-AFP：自进化智能体驱动的抗冻蛋白精准设计平台.pptx'
OUTPUT_PATH = 'EvoAgent-AFP_决赛路演PPT_郭群飞_v3.pptx'

def create_ultimate_pitch():
    """创建终极版路演PPT - 深度整合两个PPT的优势"""
    
    # 加载两个PPT
    template = Presentation(TEMPLATE_PATH)
    project = Presentation(PROJECT_PPT_PATH)
    
    print("=" * 80)
    print("🎯 开始创建终极版路演PPT")
    print("=" * 80)
    print(f"\n✅ 已加载参考模板: {len(template.slides)}页")
    print(f"✅ 已加载项目PPT: {len(project.slides)}页")
    
    # 复制模板作为基础
    shutil.copy(TEMPLATE_PATH, OUTPUT_PATH)
    prs = Presentation(OUTPUT_PATH)
    
    # ==================== 第1页: 封面 (模板第2页) ====================
    slide1 = prs.slides[1]
    
    # 从项目PPT第1页提取标题
    proj_title = ""
    for shape in project.slides[0].shapes:
        if hasattr(shape, 'text') and "EvoAgent-AFP" in shape.text:
            proj_title = shape.text.strip()
            break
    
    # 更新模板封面的标题
    for shape in slide1.shapes:
        if hasattr(shape, 'text'):
            if "参赛项目名称" in shape.text:
                shape.text = "EvoAgent-AFP"
                # 设置大标题样式
                for para in shape.text_frame.paragraphs:
                    para.font.size = Pt(48)
                    para.font.bold = True
            elif "汇报人" in shape.text or "部门" in shape.text:
                shape.text = "汇报人：郭群飞，单位：BGI华大基因华中分院"
            elif "指导老师" in shape.text:
                shape.text = "项目成员：郭群飞"
    
    print("\n✅ 第1页(封面)完成")
    
    # ==================== 第2页: 背景介绍 (模板第3页) ====================
    slide2 = prs.slides[2]
    
    # 从项目PPT第2、3页提取应用场景和痛点
    bg_content = """应用场景（三大千亿级市场）
━━━━━━━━━━━━━━━━━━━━━━━
🧬 生物医药：器官保存、细胞治疗、疫苗稳定
🍦 食品工业：冷冻食品品质保持、冷链运输  
🌾 农业育种：抗寒作物培育、畜牧精液保存

当前痛点（传统方法的瓶颈）
━━━━━━━━━━━━━━━━━━━━━━━
⏱️ 耗时漫长：实验筛选周期 3-6个月/轮
💰 成本高昂：单次突变筛选 ¥50,000+
📉 成功率低：随机突变有效率 <5%
👤 依赖专家：难以规模化、知识难传承

为什么选择AI赋能？
━━━━━━━━━━━━━━━━━━━━━━━
✓ 将数月工作压缩到数小时
✓ 研发成本降低90%以上
✓ 预测精度提升至87%
✓ 实现知识的系统化沉淀"""
    
    # 找到内容区域并更新
    for shape in slide2.shapes:
        if hasattr(shape, 'text') and len(shape.text) > 50:
            shape.text = bg_content
            for para in shape.text_frame.paragraphs:
                para.font.size = Pt(18)
                para.font.name = '微软雅黑'
                para.space_after = Pt(6)
    
    print("✅ 第2页(背景与痛点)完成")
    
    # ==================== 第3页: 解决方案 (模板第4页) ====================
    slide3 = prs.slides[3]
    
    # 从项目PPT第4页提取解决方案核心内容
    solution_content = """EvoAgent-AFP：一站式智能设计平台

核心功能亮点
━━━━━━━━━━━━━━━━━━━━━━━
🚀 多模型并行分析
   • 支持Kimi、DeepSeek、GPT等7+AI模型
   • 自动交叉验证,提升结果可靠性

⚡ 批量高效处理
   • 一次性处理500+序列
   • 并发控制1-20个任务
   • 平均单条耗时仅2-3秒

🧠 智能知识库系统
   • 分类进化知识库(AFP谱系与演化)
   • 冰结合机制库(吸附抑制理论)
   • 物理化学性质库(理化参数指标)

🔄 自进化架构
   • LLM推理引擎动态组装提示词
   • 三级递进记忆体系(短期/长期/技能)
   • 自主学习优化策略

技术突破
━━━━━━━━━━━━━━━━━━━━━━━
✓ 异步并发引擎(Semaphore控制)
✓ 智能解析器(FASTA/CSV/文本)
✓ 实时进度追踪
✓ 灵活导出(CSV/JSON)"""
    
    for shape in slide3.shapes:
        if hasattr(shape, 'text') and len(shape.text) > 30:
            if "研究内容" not in shape.text and "创新性" not in shape.text:
                shape.text = solution_content
                for para in shape.text_frame.paragraphs:
                    para.font.size = Pt(16)
                    para.font.name = '微软雅黑'
                    para.space_after = Pt(4)
    
    print("✅ 第3页(解决方案)完成")
    
    # ==================== 第4页: 创作过程 (模板第5页) ====================
    slide4 = prs.slides[4]
    
    # 从项目PPT第4页提取技术架构
    tech_content = """系统架构设计（三层验证体系）

第一层：快速初筛 (秒级)
━━━━━━━━━━━━━━━━━━━━━━━
• 几何评分引擎
• 冰结合潜力初筛与活性预测

第二层：结构验证 (分钟级)  
━━━━━━━━━━━━━━━━━━━━━━━
• PLM蛋白质语言模型
• 突变体折叠合理性与稳定性预测

第三层：机理阐释 (介观尺度)
━━━━━━━━━━━━━━━━━━━━━━━
• CGMD粗粒化分子动力学
• 自由能计算
• 在介观尺度阐释冰晶抑制机理

核心模块组成
━━━━━━━━━━━━━━━━━━━━━━━
🔧 精准突变模块
   • 识别保守功能位点与禁区
   • 每次驱动1-3个位点精确替换
   
📊 多维度量化评估
   • 禁区合规性检查
   • 冰结合面完整性判定
   • TH/IRI活性变化预测
   • 结构稳定性分析
   • 表达可行性评估
   • 安全性筛查

🧠 记忆模块（三级递进体系）
   • 短期工作记忆：记录当前轮次
   • 长期经验记忆：跨任务持久化
   • 结构化技能记忆：版本化策略

开发历程（2周快速迭代）
━━━━━━━━━━━━━━━━━━━━━━━
Week 1: 需求分析 → 架构设计 → 核心引擎开发
Week 2: 前端界面 → 系统集成 → 测试优化"""
    
    for shape in slide4.shapes:
        if hasattr(shape, 'text') and len(shape.text) > 30:
            if "两周内" not in shape.text:
                shape.text = tech_content
                for para in shape.text_frame.paragraphs:
                    para.font.size = Pt(15)
                    para.font.name = '微软雅黑'
                    para.space_after = Pt(3)
    
    print("✅ 第4页(技术架构)完成")
    
    # ==================== 第5页: 项目成果 (模板第6页) ====================
    slide5 = prs.slides[5]
    
    # 从项目PPT第5页提取对比数据
    results_content = """性能对比：传统模式 vs 智能体模式

效率提升对比
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
指标          传统方法      EvoAgent-AFP    提升倍数
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
⏱️ 分析时间    3-6个月       5分钟           ↑ 30000倍
💰 单次成本    ¥50,000      ¥50             ↓ 99.9%
📈 成功率      <5%          87%             ↑ 17倍
🔄 并发能力    1条/次       500条/次        ↑ 500倍
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

实际测试案例（100条AFP序列）
━━━━━━━━━━━━━━━━━━━━━━━
✓ 总处理时间：约3-5分钟
✓ 成功处理率：100%
✓ 平均单条耗时：2-3秒
✓ 输出质量：专业级分析报告

系统进化的四个维度
━━━━━━━━━━━━━━━━━━━━━━━
决策主体：人类驾驶 → AI驾驶
实验迭代：慢速物理 → 高速数字
知识利用：静态检索 → 动态综合
能力边界：局部扰动 → 全局探索

带来的核心价值
━━━━━━━━━━━━━━━━━━━━━━━
🎯 提高效率：研究人员从繁琐分析中解放
🎯 降低成本：大幅减少实验试错费用  
🎯 加速创新：快速验证假设,缩短研发周期
🎯 知识沉淀：可复用的分析流程和结果"""
    
    for shape in slide5.shapes:
        if hasattr(shape, 'text') and len(shape.text) > 30:
            if "实用性" not in shape.text and "落地性" not in shape.text:
                shape.text = results_content
                for para in shape.text_frame.paragraphs:
                    para.font.size = Pt(15)
                    para.font.name = '微软雅黑'
                    para.space_after = Pt(3)
    
    print("✅ 第5页(项目成果)完成")
    
    # ==================== 第6页: 其他展示 (模板第7页) ====================
    slide6 = prs.slides[6]
    
    # 从项目PPT第6、7页提取运行方式和结果
    extra_content = """运行方式（双模式支持）

🌐 Web界面模式
━━━━━━━━━━━━━━━━━━━━━━━
• 可视化操作界面
• 实时进度展示
• 交互式结果预览
• 一键导出功能

💻 CLI命令行模式
━━━━━━━━━━━━━━━━━━━━━━━
• 适合批量自动化
• 集成到工作流
• API调用支持

应用场景展望

短期应用（3-6个月）
━━━━━━━━━━━━━━━━━━━━━━━
✓ 器官移植保存液优化
✓ 细胞冷冻保护剂开发
✓ 科研工具普及

中期扩展（6-12个月）
━━━━━━━━━━━━━━━━━━━━━━━
✓ 建立AFP序列数据库
✓ 开发移动端应用
✓ 开放API接口服务
✓ 扩展到其他蛋白质类型

长期愿景（1-2年）
━━━━━━━━━━━━━━━━━━━━━━━
🚀 构建AI驱动的蛋白质研发平台
🚀 产学研深度合作生态
🚀 推动生物医药产业智能化升级

商业化路径
━━━━━━━━━━━━━━━━━━━━━━━
To B: 药企/生物科技公司授权
To R: 科研机构订阅服务
To C: 在线分析平台"""
    
    for shape in slide6.shapes:
        if hasattr(shape, 'text') and len(shape.text) > 20:
            if "其他需要展示" not in shape.text:
                shape.text = extra_content
                for para in shape.text_frame.paragraphs:
                    para.font.size = Pt(15)
                    para.font.name = '微软雅黑'
                    para.space_after = Pt(3)
    
    print("✅ 第6页(应用与展望)完成")
    
    # ==================== 保存最终版本 ====================
    prs.save(OUTPUT_PATH)
    
    print("\n" + "=" * 80)
    print("✅ 终极版PPT生成完成！")
    print("=" * 80)
    print(f"\n📁 文件位置: {OUTPUT_PATH}")
    print(f"📊 总页数: {len(prs.slides)}页")
    print(f"\n📋 PPT结构:")
    print("   第1页: 封面 - 项目名称与汇报人")
    print("   第2页: 背景介绍 - 应用场景与市场痛点")
    print("   第3页: 解决方案 - EvoAgent-AFP平台核心功能")
    print("   第4页: 技术架构 - 三层验证体系与核心模块")
    print("   第5页: 项目成果 - 性能对比与实际效果")
    print("   第6页: 应用展望 - 运行方式与商业化路径")
    print("   第7页: 补充内容")
    print("   第8页: 致谢页")
    print(f"\n✨ 本次改进亮点:")
    print("   ✓ 深度整合项目PPT的技术细节")
    print("   ✓ 保留参考模板的专业布局")
    print("   ✓ 突出核心数据和性能对比")
    print("   ✓ 逻辑清晰、层次分明")
    print(f"\n🎯 建议:")
    print("   1. 打开PPT检查排版效果")
    print("   2. 可适当添加图表和截图增强视觉效果")
    print("   3. 排练演讲控制在5分钟内")
    print("=" * 80)

if __name__ == '__main__':
    try:
        create_ultimate_pitch()
    except Exception as e:
        print(f"❌ 错误: {e}")
        import traceback
        traceback.print_exc()
